from __future__ import annotations

import statistics
from collections import Counter
from pathlib import Path
from typing import Any
from xml.etree import ElementTree
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


EMU_PER_INCH = 914400


class CoreReferenceStyleError(RuntimeError):
    pass


def resolve_core_reference(input_path: str | Path, explicit_path: str | Path | None = None) -> Path:
    """Resolve one and only one editable core-reference PPTX for a run."""
    if explicit_path:
        candidate = Path(explicit_path).expanduser().resolve()
        if not candidate.is_file():
            raise CoreReferenceStyleError(f"Core reference PPT not found: {candidate}")
        if candidate.suffix.lower() not in {".pptx", ".pptm"}:
            raise CoreReferenceStyleError("Core reference must be an editable .pptx or .pptm file.")
        return candidate

    report_dir = Path(input_path).resolve().parent
    candidates = [
        path
        for path in report_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".pptx", ".pptm"}
    ]
    preferred = [
        path
        for path in candidates
        if any(token in path.stem.lower() for token in ("核心参考", "core reference", "template", "模板"))
    ]
    if len(preferred) == 1:
        return preferred[0].resolve()
    if not preferred:
        raise CoreReferenceStyleError(
            "No core-reference PPT was found beside the report. Provide --core-reference <template.pptx>, "
            "or name exactly one candidate with 核心参考 / 模板 / template."
        )
    names = ", ".join(path.name for path in preferred)
    raise CoreReferenceStyleError(
        f"Multiple core-reference PPT candidates were found: {names}. Specify one with --core-reference."
    )


def extract_core_reference_style(reference_path: str | Path) -> dict[str, Any]:
    """Extract only the reusable visual tokens needed for effect-image generation."""
    path = Path(reference_path).resolve()
    try:
        presentation = Presentation(path)
    except Exception as exc:
        raise CoreReferenceStyleError(f"Could not read core reference PPT: {path}") from exc

    width_in = round(presentation.slide_width / EMU_PER_INCH, 3)
    height_in = round(presentation.slide_height / EMU_PER_INCH, 3)
    theme_fonts = _theme_fonts(path)
    title_records: list[dict[str, Any]] = []
    body_records: list[dict[str, Any]] = []
    caption_records: list[dict[str, Any]] = []
    colors: Counter[str] = Counter()
    fill_colors: Counter[str] = Counter()
    background_colors: Counter[str] = Counter()
    footer_boxes: list[dict[str, float]] = []
    text_block_counts: list[int] = []

    for slide in presentation.slides:
        background_rgb = _rgb_from_fill(slide.background)
        if background_rgb:
            background_colors[background_rgb] += 1
        text_blocks = 0

        for shape in slide.shapes:
            fill_rgb = _rgb_from_fill(shape)
            if fill_rgb:
                fill_colors[fill_rgb] += 1
                colors[fill_rgb] += 1

            if not getattr(shape, "has_text_frame", False):
                continue
            text = (shape.text or "").strip()
            if not text:
                continue
            text_blocks += 1
            role = _text_role(shape, presentation.slide_width, presentation.slide_height)
            record = {
                "x": round(shape.left / EMU_PER_INCH, 3),
                "y": round(shape.top / EMU_PER_INCH, 3),
                "w": round(shape.width / EMU_PER_INCH, 3),
                "h": round(shape.height / EMU_PER_INCH, 3),
                "fonts": _shape_fonts(shape, theme_fonts, role),
            }
            if role == "title":
                title_records.append(record)
            else:
                body_records.append(record)
            if shape.top >= presentation.slide_height * 0.84:
                footer_boxes.append({key: record[key] for key in ("x", "y", "w", "h")})
                caption_records.append(record)
            for font in record["fonts"]:
                if font.get("rgb"):
                    colors[font["rgb"]] += 1

        text_block_counts.append(text_blocks)

    footer_anchor = _median_box(footer_boxes)
    layout_typography = _layout_typography(path)
    palette = _palette(colors, fill_colors, background_colors)
    typography = {
        "title": _prefer_layout_typography(
            _font_summary(title_records, theme_fonts.get("major_latin") or theme_fonts.get("major_ea")),
            layout_typography.get("title"),
        ),
        "body": _font_summary(body_records, theme_fonts.get("minor_latin") or theme_fonts.get("minor_ea")),
        "caption": _font_summary(caption_records, theme_fonts.get("minor_latin") or theme_fonts.get("minor_ea")),
        "theme_major_font": theme_fonts.get("major_ea") or theme_fonts.get("major_latin"),
        "theme_minor_font": theme_fonts.get("minor_ea") or theme_fonts.get("minor_latin"),
    }
    density = {
        "median_text_blocks": round(statistics.median(text_block_counts), 1) if text_block_counts else 0,
        "recommended_evidence_modules": "2-4 on analysis pages; preserve the reference deck's intentional whitespace.",
    }
    style = {
        "template_name": path.stem,
        "reference_path": str(path),
        "slide_count": len(presentation.slides),
        "canvas": {
            "width_inches": width_in,
            "height_inches": height_in,
            "aspect_ratio": _aspect_ratio(width_in, height_in),
        },
        "palette": palette,
        "typography": typography,
        "title_system": {
            "rule": "Reuse the reference title hierarchy and conclusion-led title treatment.",
        },
        "page_chrome": {
            "rule": "Preserve the reference footer/source/page-number convention and relative placement.",
            "footer_anchor_inches": footer_anchor,
        },
        "density": density,
    }
    style["style_prompt"] = build_style_prompt(style)
    return style


def recommended_image_size(style: dict[str, Any], long_edge: int = 2048) -> str:
    """Return a gpt-image-compatible raster size matching the core PPT canvas."""
    canvas = style.get("canvas", {})
    width = float(canvas.get("width_inches") or 0)
    height = float(canvas.get("height_inches") or 0)
    if width <= 0 or height <= 0:
        return "2048x1152"
    output_height = max(16, int(round((long_edge * height / width) / 16) * 16))
    return f"{long_edge}x{output_height}"


def build_style_prompt(style: dict[str, Any]) -> str:
    canvas = style.get("canvas", {})
    palette = style.get("palette", {})
    typography = style.get("typography", {})
    title = typography.get("title", {})
    body = typography.get("body", {})
    caption = typography.get("caption", {})
    return " ".join(
        part
        for part in [
            f"Core template: {style.get('template_name', 'declared core reference') }.",
            f"Canvas: {canvas.get('aspect_ratio', 'reference aspect ratio')} ({canvas.get('width_inches', '?')} x {canvas.get('height_inches', '?')} inches).",
            f"Palette: {palette.get('prompt_summary', 'derive dominant title, support, neutral, and fill colors from the reference')}.",
            f"Typography: title {title.get('font_family', 'reference title font')} about {title.get('size_pt', 'reference size')}pt; body {body.get('font_family', 'reference body font')} about {body.get('size_pt', 'reference size')}pt; caption {caption.get('font_family', 'reference caption font')} about {caption.get('size_pt', 'reference size')}pt.",
            style.get("title_system", {}).get("rule", ""),
            style.get("page_chrome", {}).get("rule", ""),
            f"Density: {style.get('density', {}).get('recommended_evidence_modules', '')}",
        ]
        if part
    )


def _theme_fonts(path: Path) -> dict[str, str]:
    try:
        with ZipFile(path) as archive:
            root = ElementTree.fromstring(archive.read("ppt/theme/theme1.xml"))
    except Exception:
        return {}
    namespace = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    result: dict[str, str] = {}
    for role, prefix in (("major", "majorFont"), ("minor", "minorFont")):
        node = root.find(f".//a:fontScheme/a:{prefix}", namespace)
        if node is None:
            continue
        for script, tag in (("latin", "latin"), ("ea", "ea")):
            font = node.find(f"a:{tag}", namespace)
            if font is not None and font.get("typeface"):
                result[f"{role}_{script}"] = font.get("typeface")
    return result


def _layout_typography(path: Path) -> dict[str, dict[str, Any]]:
    namespace = {
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    }
    values: dict[str, list[dict[str, Any]]] = {"title": [], "body": []}
    try:
        with ZipFile(path) as archive:
            for name in archive.namelist():
                if not name.startswith("ppt/slideLayouts/slideLayout") or not name.endswith(".xml"):
                    continue
                root = ElementTree.fromstring(archive.read(name))
                for shape in root.findall(".//p:sp", namespace):
                    placeholder = shape.find(".//p:ph", namespace)
                    if placeholder is None:
                        continue
                    placeholder_type = placeholder.get("type")
                    role = "title" if placeholder_type in {"title", "ctrTitle"} else "body"
                    if placeholder_type not in {"title", "ctrTitle", "subTitle", "body"}:
                        continue
                    props = shape.find(".//a:defRPr", namespace)
                    if props is None:
                        continue
                    ea = props.find("a:ea", namespace)
                    latin = props.find("a:latin", namespace)
                    size = props.get("sz")
                    if not size and ea is None and latin is None:
                        continue
                    values[role].append(
                        {
                            "font_family": (ea.get("typeface") if ea is not None else None)
                            or (latin.get("typeface") if latin is not None else None),
                            "size_pt": round(int(size) / 100, 1) if size else None,
                        }
                    )
    except Exception:
        return {}
    result: dict[str, dict[str, Any]] = {}
    for role, records in values.items():
        names = [record["font_family"] for record in records if record.get("font_family")]
        sizes = [record["size_pt"] for record in records if record.get("size_pt")]
        if names or sizes:
            result[role] = {
                "font_family": Counter(names).most_common(1)[0][0] if names else None,
                "size_pt": round(statistics.median(sizes), 1) if sizes else None,
            }
    return result


def _text_role(shape: Any, slide_width: int, slide_height: int) -> str:
    if getattr(shape, "is_placeholder", False):
        try:
            if shape.placeholder_format.type in {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}:
                return "title"
        except Exception:
            pass
    return "title" if shape.top <= slide_height * 0.2 and shape.width >= slide_width * 0.42 else "body"


def _shape_fonts(shape: Any, theme_fonts: dict[str, str], role: str) -> list[dict[str, Any]]:
    records: list[dict[str, Any]] = []
    fallback = theme_fonts.get("major_ea") if role == "title" else theme_fonts.get("minor_ea")
    fallback = fallback or (theme_fonts.get("major_latin") if role == "title" else theme_fonts.get("minor_latin"))
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            size = round(run.font.size.pt, 1) if run.font.size else None
            rgb = _rgb_from_font(run.font)
            records.append(
                {
                    "name": run.font.name or fallback,
                    "size_pt": size,
                    "bold": bool(run.font.bold),
                    "rgb": rgb,
                }
            )
    return records or [{"name": fallback, "size_pt": None, "bold": False, "rgb": None}]


def _rgb_from_font(font: Any) -> str | None:
    try:
        rgb = font.color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def _rgb_from_fill(shape: Any) -> str | None:
    try:
        rgb = shape.fill.fore_color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def _rgb_from_line(shape: Any) -> str | None:
    try:
        rgb = shape.line.color.rgb
        return str(rgb) if rgb else None
    except Exception:
        return None


def _line_weight_pt(shape: Any) -> float | None:
    try:
        width = shape.line.width
        return round(width / 12700, 2) if width else None
    except Exception:
        return None


def _is_rounded_shape(shape: Any) -> bool:
    try:
        return "ROUNDED" in str(shape.auto_shape_type).upper()
    except Exception:
        return False


def _median_box(records: list[dict[str, Any]]) -> dict[str, float] | None:
    usable = [record for record in records if all(record.get(key) is not None for key in ("x", "y", "w", "h"))]
    if not usable:
        return None
    return {key: round(statistics.median(record[key] for record in usable), 3) for key in ("x", "y", "w", "h")}


def _shape_box(shape: Any) -> dict[str, float] | None:
    try:
        return {
            "x": round(shape.left / EMU_PER_INCH, 3),
            "y": round(shape.top / EMU_PER_INCH, 3),
            "w": round(shape.width / EMU_PER_INCH, 3),
            "h": round(shape.height / EMU_PER_INCH, 3),
        }
    except Exception:
        return None


def _is_full_slide_box(box: dict[str, float], width_in: float, height_in: float) -> bool:
    return box["w"] >= width_in * 0.98 and box["h"] >= height_in * 0.98


def _bounds_from_boxes(boxes: list[dict[str, float]]) -> dict[str, float]:
    return {
        "left": round(min(box["x"] for box in boxes), 3),
        "top": round(min(box["y"] for box in boxes), 3),
        "right": round(max(box["x"] + box["w"] for box in boxes), 3),
        "bottom": round(max(box["y"] + box["h"] for box in boxes), 3),
    }


def _median_bounds(records: list[dict[str, float]]) -> dict[str, float] | None:
    if not records:
        return None
    return {key: round(statistics.median(record[key] for record in records), 3) for key in ("left", "top", "right", "bottom")}


def _common_axis_values(values: list[float], limit: int = 6) -> list[float]:
    if not values:
        return []
    counts = Counter(round(value, 1) for value in values)
    return [value for value, _ in counts.most_common(limit)]


def _cover_composition(shapes: list[Any], width_in: float, height_in: float) -> dict[str, Any]:
    pictures: list[dict[str, float]] = []
    title_boxes: list[dict[str, float]] = []
    for shape in shapes:
        box = _shape_box(shape)
        if not box:
            continue
        if "PICTURE" in str(getattr(shape, "shape_type", "")):
            pictures.append(box)
        elif getattr(shape, "has_text_frame", False) and (shape.text or "").strip():
            if box["y"] <= height_in * 0.42 and box["w"] >= width_in * 0.2:
                title_boxes.append(box)
    if not pictures:
        return {"rule": "Use the reference cover's title hierarchy and whitespace; add an image only where the core template establishes one."}
    hero = max(pictures, key=lambda box: box["w"] * box["h"])
    hero_center = hero["x"] + hero["w"] / 2
    is_full_bleed = hero["w"] >= width_in * 0.92 and hero["h"] >= height_in * 0.92
    position = "full-bleed" if is_full_bleed else "right" if hero_center > width_in * 0.58 else "left" if hero_center < width_in * 0.42 else "center"
    title_side = "left" if title_boxes and statistics.median(box["x"] for box in title_boxes) < width_in * 0.5 else "right"
    return {
        "hero_image_anchor_inches": hero,
        "hero_image_position": position,
        "title_side": title_side,
        "rule": (
            f"Use a full-bleed editorial hero image with the title hierarchy on the {title_side}; preserve the reference crop and whitespace balance."
            if is_full_bleed
            else f"Use a {position}-anchored editorial hero image with the title hierarchy on the {title_side}; preserve the reference crop and whitespace balance."
        ),
    }


def _font_summary(records: list[dict[str, Any]], fallback: str | None) -> dict[str, Any]:
    fonts = [font for record in records for font in record.get("fonts", [])]
    names = [font["name"] for font in fonts if font.get("name")]
    sizes = [font["size_pt"] for font in fonts if font.get("size_pt")]
    return {
        "font_family": Counter(names).most_common(1)[0][0] if names else fallback,
        "size_pt": round(statistics.median(sizes), 1) if sizes else None,
        "bold_share": round(sum(1 for font in fonts if font.get("bold")) / len(fonts), 2) if fonts else 0,
    }


def _prefer_layout_typography(observed: dict[str, Any], layout: dict[str, Any] | None) -> dict[str, Any]:
    if not layout:
        return observed
    return {
        **observed,
        "font_family": layout.get("font_family") or observed.get("font_family"),
        "size_pt": layout.get("size_pt") or observed.get("size_pt"),
        "source": "core template layout" if layout.get("font_family") or layout.get("size_pt") else "slide observation",
    }


def _palette(
    colors: Counter[str],
    fills: Counter[str],
    backgrounds: Counter[str],
) -> dict[str, Any]:
    background = _top_values(backgrounds, 1, include_neutral=True)
    dominant = _top_values(colors, 8)
    fill_values = _top_values(fills, 5)
    roles = _palette_roles(dominant, fill_values, background[0] if background else "FFFFFF")
    return {
        "dominant_colors": dominant,
        "fill_colors": fill_values,
        "background_color": background[0] if background else "FFFFFF",
        "roles": roles,
        "prompt_summary": (
            f"primary {roles['primary']}; support {roles['support']}; pale fill {roles['pale_fill']}; "
            f"body text {roles['body']}; background {roles['background']}"
        ),
    }


def _palette_roles(
    dominant: list[str],
    fills: list[str],
    background: str,
) -> dict[str, str]:
    """Assign semantic color roles instead of trusting counter order."""
    all_values = _unique_hexes([*dominant, *fills])
    primary = _first_matching(all_values, _is_saturated_blue) or _first_matching(all_values, _is_dark_color) or "1B61A5"
    support = (
        _first_matching(
            [value for value in all_values if value != primary],
            lambda value: _is_saturated_blue(value) and 70 <= _luminance(value) <= 150,
        )
        or primary
    )
    pale_fill = (
        _first_matching(
            fills,
            lambda value: _luminance(value) >= 220 and _blue_bias(value) >= 8,
        )
        or _first_matching(fills, lambda value: _luminance(value) >= 220)
        or "ECF6FD"
    )
    body = (
        _first_matching(all_values, lambda value: _is_neutral(value) and _luminance(value) <= 130)
        or _first_matching(all_values, _is_dark_color)
        or "404040"
    )
    return {
        "primary": primary,
        "support": support,
        "pale_fill": pale_fill,
        "body": body,
        "background": background,
    }


def _unique_hexes(values: list[str]) -> list[str]:
    unique: list[str] = []
    for value in values:
        normalized = value.upper().removeprefix("#")
        if len(normalized) == 6 and normalized not in unique:
            unique.append(normalized)
    return unique


def _first_matching(values: list[str], predicate: Any) -> str | None:
    for value in values:
        if predicate(value):
            return value.upper().removeprefix("#")
    return None


def _rgb_tuple(value: str) -> tuple[int, int, int]:
    normalized = value.upper().removeprefix("#")
    return tuple(int(normalized[index : index + 2], 16) for index in (0, 2, 4))


def _luminance(value: str) -> float:
    red, green, blue = _rgb_tuple(value)
    return 0.2126 * red + 0.7152 * green + 0.0722 * blue


def _blue_bias(value: str) -> int:
    red, green, blue = _rgb_tuple(value)
    return blue - red + max(0, blue - green)


def _is_saturated_blue(value: str) -> bool:
    red, green, blue = _rgb_tuple(value)
    return blue >= green >= red and blue - red >= 45


def _is_neutral(value: str) -> bool:
    red, green, blue = _rgb_tuple(value)
    return max(red, green, blue) - min(red, green, blue) <= 14


def _is_dark_color(value: str) -> bool:
    return _luminance(value) <= 135


def _top_values(counter: Counter[str], limit: int, include_neutral: bool = False) -> list[str]:
    values = [value for value, _ in counter.most_common()]
    if not include_neutral:
        values = [value for value in values if value.upper() not in {"FFFFFF", "000000"}]
    return values[:limit]


def _aspect_ratio(width: float, height: float) -> str:
    if not height:
        return "unknown"
    ratio = width / height
    if abs(ratio - 16 / 9) < 0.03:
        return "16:9"
    if abs(ratio - 4 / 3) < 0.03:
        return "4:3"
    return f"{ratio:.2f}:1"


def _classify_slide(
    slide_index: int,
    has_title: bool,
    pictures: int,
    tables: int,
    charts: int,
    shapes: int,
    text_count: int,
) -> str:
    if slide_index == 1 and pictures:
        return "cover with editorial hero image"
    if tables:
        return "evidence table and commentary"
    if charts >= 2:
        return "multi-chart dashboard"
    if charts:
        return "chart and commentary"
    if shapes >= 28 and text_count >= 30:
        return "data overview with KPI modules"
    if has_title:
        return "title-led content page"
    return "section or divider page"
